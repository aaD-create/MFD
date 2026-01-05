
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from datetime import datetime
import os


SCHEMES = [
    # Equity — Parag Parikh Flexi Cap
    {
        "name": "Parag Parikh Flexi Cap Fund – Direct Growth",
        "group": "Equity", "subCat": "Flexi Cap",
        "metrics": {
            "3Y CAGR (vs category)": "23.06% vs 17.76%",
            "Sharpe": "1.55", "Beta": "0.57", "Std. Dev.": "8.41%",
            "Portfolio P/E": "20.62",
            "Exit‑load": "For units >10%: 2% ≤365d; 1% 366–730d",
            "AUM (₹ Cr)": "1,29,783", "Brokerage (bps)": "N/A",
            "Top‑3 sectors %": "48.28", "Top‑10 stocks %": "47.89",
            "Expense": "0.63%", "Turnover": "18.81%",
            "Alpha": "7.42%", "Max drawdown": "−6.05%",
            "Riskometer": "Very High"
        },
        "notes": "Core flexi‑cap with strong risk‑adjusted profile and diversified global exposure."
    },
    # Equity — HDFC Mid‑Cap
    {
        "name": "HDFC Mid‑Cap Fund – Direct Growth",
        "group": "Equity", "subCat": "Mid Cap",
        "metrics": {
            "3Y CAGR (vs category)": "26.95% vs 24.74%",
            "Sharpe": "1.316", "Beta": "0.858", "Std. Dev.": "13.847%",
            "Portfolio P/E": "—",
            "Exit‑load": "1% ≤365 days; Nil thereafter",
            "AUM (₹ Cr)": "92,168.85", "Brokerage (bps)": "N/A",
            "Top‑3 sectors %": "56.23", "Top‑10 stocks %": "33.44",
            "Expense": "1.36%", "Turnover": "18.22%",
            "Alpha": "4.77%", "Max drawdown": "−16.42%",
            "Riskometer": "Very High"
        },
        "notes": "Quality‑biased mid‑cap engine; prudent sizing for multi‑year growth."
    },
    # Hybrid — ICICI BAF (fallback if Bajaj MAF not chosen)
    {
        "name": "ICICI Prudential Balanced Advantage Fund — Regular Growth",
        "group": "Hybrid", "subCat": "Balanced Advantage",
        "metrics": {
            "3Y CAGR (vs category)": "13.61% vs 11.96%",
            "Sharpe": "1.19", "Beta": "0.55", "Std. Dev.": "5.46%",
            "Portfolio P/E": "—",
            "Exit‑load": ">30% units: 1% within 1 year",
            "AUM (₹ Cr)": "69,868", "Brokerage (bps)": "N/A",
            "Top‑3 sectors %": "41.80", "Top‑10 stocks %": "55.00",
            "Expense": "1.43% (Regular)", "Turnover": "33.97%",
            "Alpha": "3.06%", "Max drawdown": "−6.09%",
            "Riskometer": "Very High"
        },
        "notes": "Dynamic allocation moderates drawdowns; stabiliser for multi‑year horizons."
    },
    # Debt — Axis Liquid
    {
        "name": "Axis Liquid Fund – Direct Growth",
        "group": "Debt", "subCat": "Liquid",
        "metrics": {
            "3Y CAGR (vs category)": "7.08% vs 7.02%",
            "Sharpe": "3.65", "Beta": "0.38", "Std. Dev.": "~0.80%",
            "Portfolio P/E": "—",
            "Exit‑load": "Day‑1 0.0070% … Day‑7+ 0%",
            "AUM (₹ Cr)": "37,357.87", "Brokerage (bps)": "N/A",
            "Top‑3 sectors %": "—", "Top‑10 stocks %": "—",
            "Expense": "0.24% (Direct)", "Turnover": "—",
            "Alpha": "~1.30", "Max drawdown": "—",
            "Riskometer": "Low to Moderate"
        },
        "notes": "High‑liquidity sleeve to ring‑fence near‑term obligations."
    },
    # ---- Bajaj AMC schemes (enforcement targets) ----
    {
        "name": "Bajaj Finserv Multi Asset Allocation Fund — Regular Growth",
        "group": "Hybrid", "subCat": "Multi Asset",
        "metrics": {
            "3Y CAGR (vs category)": "—",
            "Sharpe": "—", "Beta": "—", "Std. Dev.": "—",
            "Portfolio P/E": "—",
            "Exit‑load": ">30% units: 1% within 1 year",
            "AUM (₹ Cr)": "1,331.01", "Brokerage (bps)": "N/A",
            "Top‑3 sectors %": "37.35", "Top‑10 stocks %": "45.36",
            "Expense": "1.99%", "Turnover": "—",
            "Alpha": "—", "Max drawdown": "—",
            "Riskometer": "Very High"
        },
        "notes": "Diversified across equity, debt & commodities (gold/silver); aims to smooth drawdowns."
    },
    {
        "name": "Bajaj Finserv Flexi Cap Fund — Regular Growth",
        "group": "Equity", "subCat": "Flexi Cap",
        "metrics": {
            "3Y CAGR (vs category)": "—",
            "Sharpe": "—", "Beta": "—", "Std. Dev.": "—",
            "Portfolio P/E": "—",
            "Exit‑load": "≤6 months: 1% (beyond 10% free units)",
            "AUM (₹ Cr)": "6,294.00", "Brokerage (bps)": "N/A",
            "Top‑3 sectors %": "—", "Top‑10 stocks %": "—",
            "Expense": "1.78%", "Turnover": "—",
            "Alpha": "—", "Max drawdown": "—",
            "Riskometer": "Very High"
        },
        "notes": "Agile core equity across market caps; long‑term megatrends via InQuBe."
    }
]

POLICY = {
    "LT1": {"Debt":100,"Hybrid":0,"Equity":0},
    "1to3": {"Debt":60,"Hybrid":40,"Equity":0},
    "3to5": {"Debt":20,"Hybrid":30,"Equity":50},
    "5to10": {"Debt":15,"Hybrid":20,"Equity":65},
    "GT10": {"Debt":10,"Hybrid":10,"Equity":80}
}

TENURE_LABELS = {"LT1":"< 1 year","1to3":"1–3 years","3to5":"3–5 years","5to10":"5–10 years","GT10":"10+ years"}


# ---------------------------
# Helper functions (no changes needed by you)
# ---------------------------
def _get_scheme(name_substring: str):
    q = name_substring.lower()
    for s in SCHEMES:
        if q in s["name"].lower():
            return s
    return None

def _select_bajaj_choice(inputs, mix) -> str:
    """
    Decide 'maf' vs 'flexi' when inputs['bajajFundChoice']=='auto' or missing.
    If equity sleeve >= hybrid sleeve -> 'flexi', else 'maf'.
    Also consider risk appetite / tolerance and liquidity needs.
    """
    choice = (inputs.get("bajajFundChoice") or "auto").lower()
    if choice in {"maf", "flexi"}:
        return choice

    eq_w = mix.get("Equity", 0)
    hy_w = mix.get("Hybrid", 0)
    risk_appetite = (inputs.get("riskAppetite") or "").lower()
    risk_tolerance = (inputs.get("riskTolerance") or "").lower()
    liquidity = (inputs.get("liquidityNeeds") or "").lower()

    if "very" in risk_appetite or "high" in risk_appetite or "aggressive" in risk_tolerance:
        return "flexi"
    if liquidity.startswith("yes") or "yes" in liquidity:
        return "maf"

    return "flexi" if eq_w >= hy_w else "maf"

def _set_tf_font(tf, size_pt=12, bold=False, italic=False, name='Calibri'):
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size_pt)
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = name


# ---------------------------
# Visual assets (optional)
# ---------------------------
MEDIA_DIR = os.path.join("static", "media")
HERO_IMG = next((os.path.join(MEDIA_DIR, f) for f in ("hero.png","hero.jpg","hero.jpeg") if os.path.exists(os.path.join(MEDIA_DIR, f))), None)
SIDE_IMG = next((os.path.join(MEDIA_DIR, f) for f in ("side.png","side.jpg","side.jpeg") if os.path.exists(os.path.join(MEDIA_DIR, f))), None)
ICON_IMG = next((os.path.join(MEDIA_DIR, f) for f in ("icon.png","icon.jpg","icon.jpeg") if os.path.exists(os.path.join(MEDIA_DIR, f))), None)

# Brand colors
ACCENT_BLUE  = RGBColor(96, 165, 250)   # #60a5fa
ACCENT_GREEN = RGBColor(34, 197, 94)    # #22c55e
HEADER_SHADE = RGBColor(230, 230, 230)  # light gray


# ---------------------------
# Main builder (this is the only function your app calls)
# ---------------------------
def build_proposal_ppt(buf, inputs):
    # --- derive asset mix ---
    tenure = inputs.get('tenure','3to5')
    mix = POLICY.get(tenure, POLICY['3to5'])

    # --- Bajaj enforcement flags from backend ---
    must_bajaj = bool(inputs.get("mustIncludeBajaj", False))
    bajaj_choice = (inputs.get("bajajFundChoice") or "auto").lower()
    if bajaj_choice == "auto":
        bajaj_choice = _select_bajaj_choice(inputs, mix)

    # --- Compose fund picks with weights ---
    picks = []
    debt_w   = mix['Debt']
    hybrid_w = mix['Hybrid']
    equity_w = mix['Equity']

    # Debt: Axis Liquid takes full debt sleeve
    if debt_w > 0:
        axis_liquid = _get_scheme("Axis Liquid Fund")
        if axis_liquid:
            picks.append((axis_liquid, debt_w))

    # Hybrid: Bajaj MAF if enforced & choice=='maf', else ICICI BAF
    if hybrid_w > 0:
        if must_bajaj and bajaj_choice == "maf":
            bajaj_maf = _get_scheme("Bajaj Finserv Multi Asset Allocation")
            if bajaj_maf:
                picks.append((bajaj_maf, hybrid_w))
        else:
            icici_baf = _get_scheme("Balanced Advantage Fund")
            if icici_baf:
                picks.append((icici_baf, hybrid_w))

    # Equity: ensure Bajaj Flexi when enforced & choice=='flexi'
    remaining = equity_w
    if equity_w > 0 and must_bajaj and bajaj_choice == "flexi":
        bajaj_flexi = _get_scheme("Bajaj Finserv Flexi Cap")
        if bajaj_flexi:
            bw = max(20, equity_w // 2)  # give a meaningful chunk
            picks.append((bajaj_flexi, bw))
            remaining -= bw
    # fill remaining equity with PPF Flexi + HDFC Mid
    if remaining > 0:
        ppf = _get_scheme("Parag Parikh Flexi Cap")
        hdfc_mid = _get_scheme("HDFC Mid‑Cap")
        w1 = remaining // 2
        w2 = remaining - w1
        if ppf:
            picks.append((ppf, w1))
        if hdfc_mid and w2 > 0:
            picks.append((hdfc_mid, w2))

    # ---------------------------
    # Build deck (Visual v3)
    # ---------------------------
    prs = Presentation()

    # Slide 1 — Cover (accent band + optional hero image)
    s1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    band = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1))
    band.fill.solid(); band.fill.fore_color.rgb = ACCENT_BLUE; band.line.fill.background()

    title_box = s1.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9.0), Inches(0.9))
    ttf = title_box.text_frame; ttf.clear()
    p = ttf.paragraphs[0]; p.text = f"Investment Proposal — {inputs.get('investorName','Investor')}"
    p.font.size = Pt(28); p.font.bold = True; p.font.name = 'Calibri'

    sub = s1.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.5), Inches(1.0))
    stf = sub.text_frame; stf.clear()
    sp = stf.paragraphs[0]
    sp.text = (
        f"Prepared by: {inputs.get('distributorName','Advisor')}\n"
        f"Date: {datetime.now().strftime('%Y-%m-%d')}\n"
        f"Risk: {inputs.get('riskAppetite','—')} / {inputs.get('riskTolerance','—')}"
    )
    _set_tf_font(stf, size_pt=14)

    if HERO_IMG:
        try:
            s1.shapes.add_picture(HERO_IMG, Inches(6.8), Inches(1.0), Inches(3.0), Inches(3.8))
        except Exception:
            pass

    # Slide 2 — Asset Mix & Suggested Allocation (bullets + optional side image)
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Asset Mix & Suggested Allocation"
    _set_tf_font(s2.shapes.title.text_frame, size_pt=20)

    tf2 = s2.placeholders[1].text_frame; tf2.clear()
    p0 = tf2.paragraphs[0]
    p0.text = f"Asset mix: Equity {equity_w}% • Hybrid {hybrid_w}% • Debt {debt_w}%"
    _set_tf_font(tf2, size_pt=12)
    for sch, wt in picks:
        pp = tf2.add_paragraph(); pp.level = 0
        pp.text = f"• {sch['name']} — {wt}% (Category: {sch['group']}/{sch['subCat']})"
    # Optional enforcement note
    if must_bajaj:
        nn = tf2.add_paragraph(); nn.level = 0
        nn.text = f"• Bajaj AMC included ({'MAF' if bajaj_choice=='maf' else 'Flexi Cap'})"
        _set_tf_font(tf2, size_pt=12)

    if SIDE_IMG:
        try:
            s2.shapes.add_picture(SIDE_IMG, Inches(7.0), Inches(1.6), Inches(2.7), Inches(3.7))
        except Exception:
            pass

    # Slides 3..N — One fund per slide (rationale + fixed table + optional icon)
    for sch, wt in picks:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        # Title
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8.5), Inches(0.6))
        tbf = tb.text_frame; tbf.clear()
        tp = tbf.paragraphs[0]; tp.text = sch["name"]
        tp.font.size = Pt(20); tp.font.bold = True; tp.font.name = 'Calibri'

        # Icon (optional)
        if ICON_IMG:
            try:
                slide.shapes.add_picture(ICON_IMG, Inches(8.7), Inches(0.2), Inches(1.2), Inches(1.2))
            except Exception:
                pass

        # Rationale beneath title (bold + italic)
        rb = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.7))
        rtf = rb.text_frame; rtf.clear()
        rp = rtf.paragraphs[0]; rp.text = sch.get('notes','Fund rationale')
        for r in rp.runs:
            r.font.bold = True; r.font.italic = True; r.font.size = Pt(12); r.font.name = 'Calibri'

        # Fixed table area
        left, top, width, height = Inches(0.5), Inches(1.8), Inches(9.0), Inches(5.0)
        metrics = sch["metrics"]
        rows = len(metrics) + 2  # + weight row
        tbl = slide.shapes.add_table(rows, 2, left, top, width, height).table
        # Header
        tbl.cell(0,0).text = "Metric"; tbl.cell(0,1).text = "Value"
        for c in (tbl.cell(0,0), tbl.cell(0,1)):
            c.fill.solid(); c.fill.fore_color.rgb = HEADER_SHADE
            _set_tf_font(c.text_frame, size_pt=12, bold=True)
        tbl.columns[0].width = Inches(3.5); tbl.columns[1].width = Inches(5.5)

        # First row: portfolio weight
        r_idx = 1
        tbl.cell(r_idx,0).text = "Portfolio weight"; tbl.cell(r_idx,1).text = f"{wt}%"
        for p in tbl.cell(r_idx,1).text_frame.paragraphs:
            p.alignment = PP_ALIGN.RIGHT
        r_idx += 1

        # Remaining metric rows (keep numeric values right-aligned)
        for k, v in metrics.items():
            tbl.cell(r_idx,0).text = k
            cv = tbl.cell(r_idx,1); cv.text = v if v is not None else '—'
            for p in cv.text_frame.paragraphs:
                p.alignment = PP_ALIGN.RIGHT
            r_idx += 1

        # Table font Calibri 12pt
        for row in tbl.rows:
            for cell in row.cells:
                _set_tf_font(cell.text_frame, size_pt=12)

        # Footer accent bar
        foot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.9), Inches(10), Inches(0.6))
        foot.fill.solid(); foot.fill.fore_color.rgb = ACCENT_GREEN; foot.line.fill.background()
        fb = slide.shapes.add_textbox(Inches(0.4), Inches(6.92), Inches(9.2), Inches(0.5))
        fbtf = fb.text_frame; fbtf.clear()
        fp = fbtf.paragraphs[0]
        fp.text = "Data as of YYYY‑MM; Sources: AMC factsheets/AMFI"
        fp.font.size = Pt(10); fp.font.name = 'Calibri'

    # Final slide — Key Metrics & Compliance (fixed box, 12pt)
    s7 = prs.slides.add_slide(prs.slide_layouts[1])
    s7.shapes.title.text = "Key Metrics & Compliance Notes"
    _set_tf_font(s7.shapes.title.text_frame, size_pt=20)
    box = s7.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(5.5))
    tf7 = box.text_frame; tf7.clear()
    bullets = [
        "Sharpe = (Return − Risk‑free) / Std. dev.; higher ⇒ better risk‑adjusted performance.",
        "Beta = sensitivity to market (1.0 = market); lower ⇒ lower volatility.",
        "Std. dev. = volatility of returns; lower ⇒ more stability.",
        "Sortino focuses on downside volatility only.",
        "Rolling returns show consistency across overlapping periods.",
        "Riskometer: SEBI/AMFI label (Low → Very High).",
        "Exit‑load: Scheme‑specific grids apply; plan redemptions to avoid charges.",
        "Disclaimer: Mutual fund investments are subject to market risks. Past performance may or may not be sustained."
    ]
    p = tf7.paragraphs[0]; p.text = "• " + bullets[0]
    for b in bullets[1:]:
        q = tf7.add_paragraph(); q.text = "• " + b
    _set_tf_font(tf7, size_pt=12)

    prs.save(buf)



